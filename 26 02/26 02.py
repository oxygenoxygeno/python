from random import randint, sample, shuffle, choices, random
from pptx import Presentation

methods = [randint, sample, shuffle, choices, random]

prs = Presentation("res.pptx")

for method in methods:
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = method.name
    subtitle.text = method.doc
prs.save("res.pptx")